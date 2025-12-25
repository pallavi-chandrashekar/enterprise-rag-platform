import io
import logging
from pathlib import Path
from typing import Iterable, Optional, Union, List, Dict
from uuid import UUID

import requests
from fastapi import status
from bs4 import BeautifulSoup  # type: ignore[import-untyped]
from sqlalchemy.orm import Session

from app.core.exceptions import AppException
from app.models.entities import Chunk, Document
from app.services.embeddings import EmbeddingService

logger = logging.getLogger(__name__)


class IngestionPipeline:
    def __init__(self, db: Session) -> None:
        self.db = db
        self.embedder = EmbeddingService()

    def process_uploaded_file(self, document: Document, file_bytes: bytes) -> None:
        try:
            extracted_content = self._extract_text(document.filename, file_bytes)
            chunks = self._chunk_content(extracted_content)
            if not chunks:
                raise ValueError("No text found in document")
            self._increment_attempt(document)
            self.process_document(document, chunks)
        except Exception as e:
            logger.error(
                f"Failed to process document {document.id} for tenant {document.tenant_id}: {e}",
                exc_info=True,
            )
            self.mark_failed(document.id, str(e))

    def process_url(self, document: Document) -> None:
        try:
            extracted_content = self._extract_text(document.filename)
            chunks = self._chunk_content(extracted_content)
            if not chunks:
                raise ValueError("No text found in document")
            self._increment_attempt(document)
            self.process_document(document, chunks)
        except Exception as e:
            logger.error(
                f"Failed to process document {document.id} for tenant {document.tenant_id}: {e}",
                exc_info=True,
            )
            self.mark_failed(document.id, str(e))

    def process_document(self, document: Document, chunks: list[str]) -> None:
        embeddings = self.embedder.embed_texts(chunks)
        with self.db.begin():
            # If retrying, clear prior chunks for this document to avoid duplicates.
            self.db.query(Chunk).filter(Chunk.document_id == document.id, Chunk.tenant_id == document.tenant_id).delete(synchronize_session=False)
            chunk_records = []
            for content, embedding in zip(chunks, embeddings):
                chunk_records.append(
                    Chunk(
                        tenant_id=document.tenant_id,
                        kb_id=document.kb_id,
                        document_id=document.id,
                        content=content,
                        embedding=embedding,
                    )
                )
            self.db.add_all(chunk_records)
            document.status = "READY"
            self._merge_metadata(document, {"last_error": None})
            self.db.add(document)

    def mark_failed(self, document_id: UUID, reason: str) -> None:
        doc = self.db.get(Document, document_id)
        if not doc:
            return
        with self.db.begin():
            doc.status = "FAILED"
            self._merge_metadata(doc, {"last_error": reason})
            self.db.add(doc)

    def _extract_text(self, filename: str, data: Optional[bytes] = None) -> Union[str, List[Dict[str, str]]]:
        if filename.startswith("http"):
            try:
                response = requests.get(filename, timeout=15)
                response.raise_for_status()
                data = response.content # Assign response content to data
                # Now process data as HTML below
            except requests.RequestException as e:
                raise ValueError(f"Failed to fetch URL {filename}: {e}") from e

        if data is None:
            raise ValueError("File content is missing for non-URL ingestion.")

        ext = Path(filename).suffix.lower()
        if ext in {".txt", ".md", ".text"}:
            return data.decode("utf-8", errors="ignore")
        if ext == ".pdf":
            try:
                from pypdf import PdfReader
            except ImportError as exc:  # pragma: no cover - runtime guard
                raise AppException(detail="pypdf not installed", status_code=status.HTTP_500_INTERNAL_SERVER_ERROR) from exc
            reader = PdfReader(io.BytesIO(data))
            return "\n".join(page.extract_text() or "" for page in reader.pages)
        if ext == ".docx":
            try:
                import docx
            except ImportError as exc:  # pragma: no cover - runtime guard
                raise AppException(detail="python-docx not installed", status_code=status.HTTP_500_INTERNAL_SERVER_ERROR) from exc
            doc = docx.Document(io.BytesIO(data))
            return "\n".join(p.text for p in doc.paragraphs)
        if ext == ".pptx":
            try:
                import pptx
            except ImportError as exc:  # pragma: no cover - runtime guard
                raise AppException(detail="python-pptx not installed", status_code=status.HTTP_500_INTERNAL_SERVER_ERROR) from exc
            prs = pptx.Presentation(io.BytesIO(data))
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_runs.append(run.text)
            return "\n".join(text_runs)
        if ext in {".html", ".htm"} or filename.startswith("http"): # Handle URL content as HTML
            soup = BeautifulSoup(data, "html.parser")
            for tag in soup(["script", "style", "head", "title", "meta", "link"]):
                tag.decompose()
            
            structured_content: List[Dict[str, str]] = []
            for tag_name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'li']:
                for tag in soup.find_all(tag_name):
                    text = tag.get_text(separator=" ", strip=True)
                    if text:
                        structured_content.append({"type": tag_name, "content": text})
            return structured_content
        raise ValueError(f"Unsupported file type: {ext or 'unknown'}")

    def _chunk_content(self, content: Union[str, List[Dict[str, str]]], max_words: int = 220, overlap: int = 40) -> list[str]:
        if isinstance(content, str):
            # Existing word-based chunking for flat strings
            words = content.split()
            if not words:
                return []
            chunks: list[str] = []
            start = 0
            step = max_words - overlap if max_words > overlap else max_words
            while start < len(words):
                end = start + max_words
                slice_words = words[start:end]
                if not slice_words:
                    break
                chunk = " ".join(slice_words).strip()
                if chunk:
                    chunks.append(chunk)
                if len(slice_words) < overlap and start > 0:
                    break
                start += step
            return chunks
        elif isinstance(content, list):
            # Hierarchy-aware chunking for structured content
            chunks: list[str] = []
            current_chunk_words: list[str] = []
            current_heading = ""

            for block in content:
                block_type = block['type']
                block_content = block['content']
                block_words = block_content.split()

                if block_type.startswith('h'):
                    current_heading = block_content
                    # Start a new chunk if current_chunk_words is not empty
                    if current_chunk_words:
                        chunks.append(" ".join(current_chunk_words).strip())
                        current_chunk_words = []
                    # Add heading itself as a chunk if it's long enough
                    if len(block_words) > 5: # Arbitrary threshold for a meaningful heading chunk
                        chunks.append(block_content)
                    else: # If short, prepend to next content
                        current_chunk_words.extend(block_words)
                elif block_type in ['p', 'li']:
                    # Prepend current heading to paragraph content for context
                    prefixed_block = f"{current_heading}: {block_content}" if current_heading else block_content
                    prefixed_block_words = prefixed_block.split()
                    
                    if len(current_chunk_words) + len(prefixed_block_words) <= max_words:
                        current_chunk_words.extend(prefixed_block_words)
                    else:
                        if current_chunk_words:
                            chunks.append(" ".join(current_chunk_words).strip())
                        current_chunk_words = prefixed_block_words
                        
                        # Handle long individual blocks
                        while len(current_chunk_words) > max_words:
                            chunks.append(" ".join(current_chunk_words[:max_words]).strip())
                            current_chunk_words = current_chunk_words[max_words - overlap:] # Add overlap

            if current_chunk_words:
                chunks.append(" ".join(current_chunk_words).strip())

            # Filter out empty chunks and ensure minimal length
            return [chunk for chunk in chunks if chunk and len(chunk.split()) >= 5]
        return []

    def _increment_attempt(self, document: Document) -> None:
        attempts = 0
        if document.doc_metadata and "ingestion_attempts" in document.doc_metadata:
            try:
                attempts = int(document.doc_metadata["ingestion_attempts"])
            except (TypeError, ValueError):
                attempts = 0
        attempts += 1
        self._merge_metadata(document, {"ingestion_attempts": attempts})
        document.status = "PROCESSING"
        self.db.add(document)

    def _merge_metadata(self, document: Document, updates: dict[str, object]) -> None:
        meta = document.doc_metadata.copy() if document.doc_metadata else {}
        meta.update(updates)
        document.doc_metadata = meta
