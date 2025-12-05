import uuid
import io
import pptx
from unittest.mock import patch

import pytest
from fastapi import HTTPException

from app.api.routes import _parse_metadata
from app.services.ingestion import IngestionPipeline


class DummyDB:
    """Minimal stub to satisfy pipeline constructor for unit-only tests."""


def test_chunk_text_overlaps():
    text = " ".join(str(i) for i in range(1, 26))  # 25 words
    pipeline = IngestionPipeline(db=DummyDB())

    chunks = pipeline._chunk_text(text, max_words=10, overlap=2)

    assert len(chunks) == 4
    # Overlap should preserve some words across adjacent chunks.
    assert chunks[0].split()[-2:] == chunks[1].split()[:2]


def test_extract_text_plain_text():
    pipeline = IngestionPipeline(db=DummyDB())
    content = "Hello world\nthis is a test."
    extracted = pipeline._extract_text("sample.txt", content.encode("utf-8"))
    assert "Hello world" in extracted
    assert "this is a test." in extracted


def test_extract_text_html_strips_scripts():
    pipeline = IngestionPipeline(db=DummyDB())
    html = b"""
    <html><head><script>var bad = 1;</script></head>
    <body><p>Hello</p><p>world</p></body></html>
    """
    extracted = pipeline._extract_text("page.html", html)
    assert "Hello world" in extracted
    assert "bad" not in extracted


def test_extract_text_pptx():
    pipeline = IngestionPipeline(db=DummyDB())
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    box = slide.shapes.add_textbox(1, 2, 3, 4)
    box.text_frame.text = "This is a test presentation."
    
    f = io.BytesIO()
    prs.save(f)
    f.seek(0)
    
    extracted = pipeline._extract_text("sample.pptx", f.read())
    assert "This is a test presentation." in extracted


@patch("requests.get")
def test_extract_text_from_url(mock_get):
    pipeline = IngestionPipeline(db=DummyDB())
    url = "https://example.com"
    html_content = "<html><body><p>This is a test from a URL.</p></body></html>"
    mock_get.return_value.status_code = 200
    mock_get.return_value.content = html_content.encode("utf-8")

    extracted = pipeline._extract_text(url)
    assert "This is a test from a URL." in extracted


def test_chunk_text_allows_short_docs():
    pipeline = IngestionPipeline(db=DummyDB())
    text = "short doc with only a few words"
    chunks = pipeline._chunk_text(text, max_words=220, overlap=40)
    assert len(chunks) == 1


def test_parse_metadata_valid_and_invalid():
    assert _parse_metadata('{"a":1}') == {"a": 1}
    assert _parse_metadata("") is None
    with pytest.raises(HTTPException):
        _parse_metadata("not-json")
